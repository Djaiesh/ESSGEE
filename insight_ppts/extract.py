"""Extract all text content from PPTX files, slide by slide."""
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_DIR = r"D:\new_es\ESSGEE\insight_ppts"
FILES = [
    "20260705 Article 1_Bridging Governance, Strategy & Delivery R1.pptx",
    "20260705 Article 2_Strategy Through Sustainability.pptx",
    "20260705 Article 3_Systems & Compliance.pptx",
]

def extract_text_from_shape(shape, indent=0):
    """Recursively extract text from a shape, including grouped shapes and tables."""
    texts = []
    prefix = "  " * indent

    # Group shapes: recurse into children
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            texts.extend(extract_text_from_shape(child, indent))
        return texts

    # Tables
    if shape.has_table:
        table = shape.table
        for row_idx, row in enumerate(table.rows):
            row_texts = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_texts.append(cell_text)
            if row_texts:
                texts.append(f"{prefix}[Table Row {row_idx}] " + " | ".join(row_texts))
        return texts

    # Text frames (most shapes)
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            para_text = para.text.strip()
            if para_text:
                texts.append(f"{prefix}{para_text}")

    return texts


def extract_pptx(filepath):
    """Extract all text from a PPTX file, returning a list of (slide_num, texts) tuples."""
    prs = Presentation(filepath)
    slides_data = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            slide_texts.extend(extract_text_from_shape(shape))
        slides_data.append((slide_num, slide_texts))
    return slides_data


def main():
    separator = "=" * 100
    for fname in FILES:
        fpath = os.path.join(PPTX_DIR, fname)
        print(separator)
        print(f"FILE: {fname}")
        print(separator)
        if not os.path.isfile(fpath):
            print(f"  *** FILE NOT FOUND: {fpath} ***\n")
            continue

        slides = extract_pptx(fpath)
        for slide_num, texts in slides:
            print(f"\n--- Slide {slide_num} ---")
            if texts:
                for t in texts:
                    print(t)
            else:
                print("  (no text content)")
        print("\n")


if __name__ == "__main__":
    main()
